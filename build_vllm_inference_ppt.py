from __future__ import annotations

from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
import seaborn as sns
from matplotlib import patches
from matplotlib.lines import Line2D
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("vLLM框架推理流程及其性能增益可视化分析.pptx")
ASSET_DIR = Path("vllm_ppt_assets")

SLATE = "#708090"
ROYAL = "#4169E1"
EMERALD = "#50C878"
LIGHT_BG = "#F7F9FB"
PALE_RED = "#F4C7C3"
GRID = "#D7DEE8"
TEXT = "#1F2933"

RGB = {
    "slate": RGBColor(0x70, 0x80, 0x90),
    "royal": RGBColor(0x41, 0x69, 0xE1),
    "emerald": RGBColor(0x50, 0xC8, 0x78),
    "text": RGBColor(0x1F, 0x29, 0x33),
    "muted": RGBColor(0x5A, 0x66, 0x73),
    "light": RGBColor(0xF7, 0xF9, 0xFB),
}

plt.rcParams.update(
    {
        "font.family": ["Arial", "DejaVu Sans"],
        "axes.edgecolor": SLATE,
        "axes.labelcolor": TEXT,
        "axes.titlesize": 15,
        "axes.titleweight": "bold",
        "axes.labelsize": 11,
        "xtick.color": TEXT,
        "ytick.color": TEXT,
        "legend.frameon": True,
        "legend.framealpha": 0.96,
        "figure.facecolor": "white",
        "savefig.facecolor": "white",
    }
)
sns.set_theme(style="whitegrid")


def save_fig(fig: plt.Figure, name: str, tight: bool = True) -> Path:
    ASSET_DIR.mkdir(exist_ok=True)
    path = ASSET_DIR / f"{name}.png"
    if tight:
        fig.savefig(path, dpi=260, bbox_inches="tight")
    else:
        fig.subplots_adjust(left=0.025, right=0.985, top=0.89, bottom=0.055)
        fig.savefig(path, dpi=260)
    plt.close(fig)
    return path


def add_rect(
    ax,
    xy,
    width,
    height,
    color,
    label=None,
    text=None,
    hatch=None,
    alpha=1.0,
    lw=1.2,
    ec="white",
    fontsize=9,
    text_color=TEXT,
):
    rect = patches.Rectangle(
        xy,
        width,
        height,
        linewidth=lw,
        edgecolor=ec,
        facecolor=color,
        hatch=hatch,
        alpha=alpha,
        label=label,
    )
    ax.add_patch(rect)
    if text:
        ax.text(
            xy[0] + width / 2,
            xy[1] + height / 2,
            text,
            ha="center",
            va="center",
            fontsize=fontsize,
            color=text_color,
            weight="bold" if fontsize >= 9 else "normal",
        )
    return rect


def figure_title(ax, title: str, subtitle: str | None = None) -> None:
    ax.set_title(title, loc="left", color=TEXT, pad=14)
    if subtitle:
        ax.text(
            0,
            1.02,
            subtitle,
            transform=ax.transAxes,
            fontsize=10,
            color=SLATE,
            va="bottom",
        )


def build_slide_1_figure() -> Path:
    fig, ax = plt.subplots(figsize=(10.6, 6.4))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 7)
    ax.axis("off")
    figure_title(ax, "Paged Memory Topology", "Logical continuity is preserved while physical KV blocks are scattered.")

    logical_positions = [(1.0, 5.5 - i * 0.72) for i in range(6)]
    physical_positions = {
        0: (6.6, 5.2),
        1: (7.8, 3.0),
        2: (5.8, 2.1),
        3: (8.3, 5.7),
        4: (6.2, 3.9),
        5: (8.7, 2.0),
    }
    colors = [ROYAL, EMERALD, SLATE, ROYAL, EMERALD, SLATE]
    mapping = {0: 3, 1: 0, 2: 4, 3: 1, 4: 5, 5: 2}

    ax.text(1.35, 6.35, "Logical Blocks", ha="center", fontsize=12, weight="bold", color=TEXT)
    ax.text(7.25, 6.35, "Physical Blocks in GPU KV Pool", ha="center", fontsize=12, weight="bold", color=TEXT)
    ax.text(4.35, 0.75, "Block Table: logical block id -> physical block id", ha="center", fontsize=10, color=SLATE)

    for i, (x, y) in enumerate(logical_positions):
        add_rect(ax, (x, y), 0.72, 0.46, colors[i], text=f"L{i}", fontsize=10)
        ax.text(x + 0.9, y + 0.23, f"seq[{i * 16}:{(i + 1) * 16})", va="center", fontsize=8.8, color=SLATE)

    for pid, (x, y) in physical_positions.items():
        add_rect(ax, (x, y), 0.78, 0.54, colors[pid % 3], text=f"P{pid}", fontsize=10)
        ax.text(x + 0.39, y - 0.16, f"addr 0x{0xA00 + pid * 0x80:X}", ha="center", va="top", fontsize=7.7, color=SLATE)

    for logical_id, physical_id in mapping.items():
        lx, ly = logical_positions[logical_id]
        px, py = physical_positions[physical_id]
        ax.annotate(
            "",
            xy=(px, py + 0.28),
            xytext=(lx + 0.72, ly + 0.23),
            arrowprops=dict(arrowstyle="->", lw=1.1, color=colors[logical_id], alpha=0.82),
        )
        mx = (lx + px) / 2 + 0.15
        my = (ly + py) / 2 + 0.05
        ax.text(mx, my, f"L{logical_id}->P{physical_id}", fontsize=7.4, color=TEXT)

    for row in range(3):
        for col in range(4):
            x = 5.6 + col * 1.05
            y = 1.35 + row * 1.55
            ax.add_patch(
                patches.Rectangle((x, y), 0.92, 0.78, facecolor="none", edgecolor=GRID, linewidth=0.8)
            )

    legend = [
        Line2D([0], [0], color=ROYAL, lw=5, label="Request A blocks"),
        Line2D([0], [0], color=EMERALD, lw=5, label="Request B blocks"),
        Line2D([0], [0], color=SLATE, lw=5, label="Shared/free-list blocks"),
    ]
    ax.legend(handles=legend, loc="lower right", fontsize=8.7)
    return save_fig(fig, "slide_01_paged_memory_topology")


def build_slide_2_figure() -> Path:
    fig, ax = plt.subplots(figsize=(10.7, 6.4))
    ax.set_xlim(0, 100)
    ax.set_ylim(0, 4.2)
    ax.set_xlabel("GPU memory address space (GB, normalized)")
    ax.set_yticks([3.05, 1.25])
    ax.set_yticklabels(["Traditional pre-allocation", "vLLM block allocator"])
    figure_title(ax, "KV Cache Fragmentation: Micro-level Comparison")
    ax.grid(axis="x", color=GRID, linewidth=0.8)
    ax.grid(axis="y", visible=False)

    traditional = [
        (0, 10, ROYAL, "Live KV\n10GB", None),
        (10, 13, "#B9C3CF", "Internal\n13GB", None),
        (23, 8, EMERALD, "Live KV\n8GB", None),
        (31, 10, "#DCE2EA", "External\n10GB", "///"),
        (41, 12, ROYAL, "Live KV\n12GB", None),
        (53, 39, PALE_RED, "Reserved but unused\n39GB", None),
        (92, 8, "#DCE2EA", "External\n8GB", "///"),
    ]
    for x, w, color, label, hatch in traditional:
        add_rect(ax, (x, 2.55), w, 0.9, color, text=label, hatch=hatch, fontsize=8.1, text_color=TEXT, ec="white")
        ax.text(x + w / 2, 3.58, f"{w}%", ha="center", fontsize=8, color=TEXT)

    block_w = 3.0
    x = 0
    block_id = 0
    while x < 96:
        color = [ROYAL, EMERALD, SLATE][block_id % 3]
        add_rect(ax, (x, 0.82), min(block_w, 96 - x), 0.86, color, text=f"B{block_id}" if block_id < 8 else None, fontsize=6.8)
        x += block_w
        block_id += 1
    add_rect(ax, (96, 0.82), 4, 0.86, "#EEF2F6", text="free\n4%", fontsize=8, ec=SLATE, text_color=SLATE)

    ax.text(3, 3.85, "Utilization: 20%-40% (example: 30%)", color=SLATE, fontsize=10.5, weight="bold")
    ax.text(3, 2.02, "Utilization: >96% with near-contiguous physical block packing", color=EMERALD, fontsize=10.5, weight="bold")

    handles = [
        patches.Patch(facecolor=ROYAL, label="Active KV cache"),
        patches.Patch(facecolor="#B9C3CF", label="Internal fragmentation"),
        patches.Patch(facecolor="#DCE2EA", hatch="///", label="External fragmentation"),
        patches.Patch(facecolor=PALE_RED, label="Reserved / unused"),
        patches.Patch(facecolor=EMERALD, label="Paged physical blocks"),
    ]
    ax.legend(handles=handles, loc="upper right", fontsize=8.3, ncol=2)
    for spine in ["top", "right", "left"]:
        ax.spines[spine].set_visible(False)
    return save_fig(fig, "slide_02_kv_fragmentation")


def build_slide_3_figure() -> Path:
    fig, ax = plt.subplots(figsize=(9.55, 7.0))
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 8)
    ax.axis("off")
    figure_title(ax, "Prefill Phase: Tensor Projection and Paged Block Allocation")

    ax.text(1.65, 7.05, "1. Prompt tensor", ha="center", fontsize=11.5, weight="bold", color=TEXT)
    rows, cols = 5, 10
    for i in range(rows):
        for j in range(cols):
            color = [ROYAL, EMERALD, SLATE][(i + j) % 3]
            add_rect(ax, (0.35 + j * 0.24, 5.92 - i * 0.24), 0.21, 0.21, color, alpha=0.78, lw=0.25)
    ax.text(1.55, 5.18, "X_prompt: [B=1, S=96, d_model=4096]", ha="center", fontsize=8.4, color=SLATE)
    for j, token in enumerate(["t0", "t15", "t16", "t31", "...", "t95"]):
        ax.text(0.45 + j * 0.43, 4.9, token, fontsize=6.8, color=SLATE)

    add_rect(ax, (3.05, 5.35), 1.4, 1.35, "#EEF2F6", text="Embedding\n+ position", fontsize=9.4, ec=SLATE, text_color=TEXT)
    ax.annotate("", xy=(3.05, 6.02), xytext=(2.83, 6.02), arrowprops=dict(arrowstyle="->", lw=1.3, color=SLATE))
    add_rect(ax, (4.85, 5.35), 1.55, 1.35, "#EEF2F6", text="Transformer\nlayers", fontsize=9.4, ec=SLATE, text_color=TEXT)
    ax.annotate("", xy=(4.85, 6.02), xytext=(4.45, 6.02), arrowprops=dict(arrowstyle="->", lw=1.3, color=SLATE))
    add_rect(ax, (6.85, 5.35), 1.6, 1.35, "#EEF2F6", text="Q,K,V\nprojection", fontsize=9.4, ec=SLATE, text_color=TEXT)
    ax.text(7.65, 5.07, "W_q, W_k, W_v", ha="center", fontsize=7.8, color=SLATE)
    ax.annotate("", xy=(6.85, 6.02), xytext=(6.4, 6.02), arrowprops=dict(arrowstyle="->", lw=1.3, color=SLATE))

    ax.text(10.4, 7.05, "2. KV cache tensor", ha="center", fontsize=11.5, weight="bold", color=TEXT)
    for layer in range(4):
        x = 9.0 + layer * 0.25
        y = 5.25 + layer * 0.12
        add_rect(ax, (x, y), 2.9, 0.72, ROYAL if layer % 2 == 0 else EMERALD, alpha=0.78, lw=0.7)
    ax.text(10.7, 5.67, "K/V: [n_layers=32, H_kv=32,\nS=96, d_head=128]", ha="center", va="center", fontsize=8.2, color="white", weight="bold")
    ax.annotate("", xy=(9.0, 5.94), xytext=(8.45, 5.94), arrowprops=dict(arrowstyle="->", lw=1.3, color=SLATE))

    ax.text(2.9, 4.25, "3. Logical slicing (Block Size = 16 tokens)", fontsize=11.2, weight="bold", color=TEXT)
    logical_blocks = []
    phys_ids = [12, 29, 4, 41, 8, 33]
    for b in range(6):
        x = 0.62 + b * 1.02
        logical_blocks.append((x, 3.42))
        add_rect(ax, (x, 3.42), 0.86, 0.56, [ROYAL, EMERALD, SLATE][b % 3], text=f"L{b}", fontsize=8.8)
        ax.text(x + 0.43, 3.23, f"{b*16}-{b*16+15}", ha="center", fontsize=7.0, color=SLATE)
    ax.text(0.65, 3.02, "logical_token_id -> logical_block_id", fontsize=7.8, color=SLATE)

    ax.text(8.0, 4.25, "4. Block table materialized during prefill", fontsize=11.2, weight="bold", color=TEXT)
    add_rect(ax, (7.1, 2.0), 2.7, 1.95, "white", ec=SLATE, lw=1.1)
    add_rect(ax, (7.1, 3.55), 2.7, 0.4, SLATE, text="logical id | physical page | base addr", fontsize=7.1)
    for b, pid in enumerate(phys_ids):
        y = 3.18 - b * 0.25
        ax.text(7.28, y, f"L{b}", fontsize=7.2, color=TEXT)
        ax.text(8.05, y, f"P{pid}", fontsize=7.2, color=TEXT)
        ax.text(8.82, y, f"0x{0xA000 + pid * 0x100:X}", fontsize=7.2, color=TEXT)
        if b < 6:
            ax.plot([logical_blocks[b][0] + 0.43, 7.1], [3.42, y + 0.03], color=[ROYAL, EMERALD, SLATE][b % 3], lw=0.75, alpha=0.7)

    ax.text(12.0, 4.25, "5. GPU KV pool", fontsize=11.2, weight="bold", color=TEXT, ha="center")
    pool_positions = {}
    pid_grid = [2, 12, 19, 29, 4, 55, 41, 6, 8, 21, 33, 48]
    for r in range(3):
        for c in range(4):
            idx = r * 4 + c
            pid = pid_grid[idx]
            x = 10.4 + c * 0.72
            y = 3.22 - r * 0.55
            pool_positions[pid] = (x, y)
            is_used = pid in phys_ids
            color = [ROYAL, EMERALD, SLATE][phys_ids.index(pid) % 3] if is_used else "#EEF2F6"
            label = f"P{pid}" if is_used else "free"
            add_rect(ax, (x, y), 0.62, 0.42, color, text=label, fontsize=6.8, ec=GRID if not is_used else "white", text_color=TEXT if not is_used else "white")
    for b, pid in enumerate(phys_ids):
        if pid in pool_positions:
            x, y = pool_positions[pid]
            ax.annotate("", xy=(x, y + 0.22), xytext=(9.8, 3.18 - b * 0.25), arrowprops=dict(arrowstyle="->", lw=0.85, color=[ROYAL, EMERALD, SLATE][b % 3], alpha=0.75))

    ax.text(
        7.0,
        0.95,
        "logical_block_id = floor(token_id / 16)  |  physical_addr = block_table[logical_block_id].base + slot_id * bytes(KV_token)",
        ha="center",
        fontsize=8.7,
        color=TEXT,
        bbox=dict(boxstyle="round,pad=0.35", fc=LIGHT_BG, ec=GRID),
    )
    ax.text(0.55, 0.52, "Allocator state after prefill: 6 blocks allocated, no max-length reservation, free-list remains reusable.", fontsize=8.4, color=SLATE)
    return save_fig(fig, "slide_03_prefill_blocks", tight=False)


def build_slide_4_figure() -> Path:
    fig, ax = plt.subplots(figsize=(9.55, 7.0))
    ax.set_xlim(0, 18)
    ax.set_ylim(0, 8)
    ax.axis("off")
    figure_title(ax, "Decoding Phase: Slot-level Fill and Allocator State")

    ax.text(2.2, 7.15, "Active request state", ha="center", fontsize=11.5, weight="bold", color=TEXT)
    add_rect(ax, (0.45, 5.63), 3.55, 1.05, "white", ec=SLATE, lw=1.1)
    add_rect(ax, (0.45, 6.28), 3.55, 0.4, SLATE, text="req | len | last block | fill", fontsize=7.6)
    req_rows = [("R_A", "63", "P42", "15/16"), ("R_B", "28", "P17", "12/16")]
    for i, row in enumerate(req_rows):
        y = 6.0 - i * 0.32
        ax.text(0.68, y, row[0], fontsize=7.7, color=TEXT)
        ax.text(1.35, y, row[1], fontsize=7.7, color=TEXT)
        ax.text(2.02, y, row[2], fontsize=7.7, color=TEXT)
        ax.text(2.92, y, row[3], fontsize=7.7, color=TEXT)

    ax.text(7.2, 7.15, "Physical Block P42: 16 slots", ha="center", fontsize=12.3, weight="bold", color=TEXT)
    start_x, start_y = 4.1, 5.25
    slot_w, slot_h = 0.62, 0.52
    for i in range(16):
        row = 0 if i < 8 else 1
        col = i if i < 8 else i - 8
        x = start_x + col * slot_w
        y = start_y - row * 0.62
        color = ROYAL if i < 15 else EMERALD
        alpha = 0.93 if i < 15 else 0.6
        add_rect(ax, (x, y), slot_w - 0.04, slot_h, color, text=str(i), fontsize=7.7, alpha=alpha)
    ax.text(6.53, 4.26, "slot 15 receives K_t/V_t", fontsize=8.2, color=EMERALD, weight="bold")

    add_rect(ax, (10.15, 5.35), 1.35, 0.82, "#EEF2F6", text="x_t\nnew token", fontsize=8.6, ec=SLATE, text_color=TEXT)
    ax.annotate("", xy=(12.0, 5.76), xytext=(11.5, 5.76), arrowprops=dict(arrowstyle="->", lw=1.2, color=SLATE))
    add_rect(ax, (12.0, 5.18), 2.15, 1.16, EMERALD, text="K_t, V_t\n[2,H_kv,d_head]", fontsize=8.6)
    ax.annotate("", xy=(start_x + 7 * slot_w + 0.31, start_y - 0.62), xytext=(12.0, 5.18), arrowprops=dict(arrowstyle="->", lw=1.8, color=EMERALD))

    ax.text(2.2, 4.42, "Block table lookup", ha="center", fontsize=11.2, weight="bold", color=TEXT)
    add_rect(ax, (0.55, 2.82), 3.35, 1.22, "white", ec=SLATE, lw=1.1)
    add_rect(ax, (0.55, 3.68), 3.35, 0.36, SLATE, text="logical block | physical | fill pointer", fontsize=7.2)
    for i, row in enumerate([("L1", "P29", "full"), ("L2", "P42", "slot=15")]):
        y = 3.4 - i * 0.36
        ax.text(0.88, y, row[0], fontsize=7.7, color=TEXT)
        ax.text(1.83, y, row[1], fontsize=7.7, color=TEXT)
        ax.text(2.73, y, row[2], fontsize=7.7, color=TEXT)
    ax.annotate("", xy=(4.1, 4.78), xytext=(3.9, 3.06), arrowprops=dict(arrowstyle="->", lw=1.2, color=ROYAL))

    ax.text(7.3, 3.8, "Address arithmetic", ha="center", fontsize=11.2, weight="bold", color=TEXT)
    formulas = [
        "token_id = 63, block_size = 16",
        "logical_block = floor(63/16) = L3",
        "slot_id = 63 mod 16 = 15",
        "addr = base(P42) + 15 * bytes(KV_token)",
    ]
    for i, txt in enumerate(formulas):
        ax.text(5.0, 3.38 - i * 0.34, txt, fontsize=8.3, color=TEXT)
    ax.text(5.0, 1.78, "bytes(KV_token)=2*n_layers*H_kv*d_head*dtype_bytes", fontsize=8.1, color=SLATE)

    ax.text(13.4, 3.8, "Allocator transition", ha="center", fontsize=11.2, weight="bold", color=TEXT)
    add_rect(ax, (11.45, 2.9), 1.15, 0.56, ROYAL, text="P42\nfull", fontsize=7.4)
    ax.annotate("", xy=(13.55, 3.18), xytext=(12.6, 3.18), arrowprops=dict(arrowstyle="->", lw=1.2, color=SLATE))
    add_rect(ax, (13.55, 2.9), 1.15, 0.56, EMERALD, text="P77\nnew", fontsize=7.4)
    ax.text(12.85, 2.48, "next token triggers\nfree-list pop", ha="center", fontsize=7.7, color=SLATE)

    timeline = [("sample x_t", ROYAL), ("project K/V", EMERALD), ("write slot", EMERALD), ("attention", SLATE), ("sample x_{t+1}", ROYAL)]
    for i, (label, color) in enumerate(timeline):
        x = 2.0 + i * 2.55
        add_rect(ax, (x, 0.78), 1.55, 0.42, color, text=label, fontsize=7.2)
        if i < len(timeline) - 1:
            ax.annotate("", xy=(x + 2.42, 0.99), xytext=(x + 1.55, 0.99), arrowprops=dict(arrowstyle="->", lw=1.0, color=SLATE))

    handles = [
        patches.Patch(facecolor=ROYAL, label="Existing KV slots / previous steps"),
        patches.Patch(facecolor=EMERALD, alpha=0.65, label="Current step insertion and next page"),
        patches.Patch(facecolor="#EEF2F6", edgecolor=SLATE, label="Scheduler metadata"),
    ]
    ax.legend(handles=handles, loc="lower left", fontsize=8.2)
    return save_fig(fig, "slide_04_decoding_slot_fill", tight=False)


def build_slide_5_figure() -> Path:
    fig, ax = plt.subplots(figsize=(9.55, 7.0))
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 8)
    ax.axis("off")
    figure_title(ax, "PagedAttention Lookup Logic: From Logical Tokens to Gathered KV")

    ax.text(1.8, 7.1, "Logical sequence R_i", fontsize=11.5, weight="bold", color=TEXT, ha="center")
    logical = []
    phys = [7, 2, 11, 4, 14, 9]
    for i in range(6):
        y = 6.28 - i * 0.56
        logical.append((0.55, y))
        add_rect(ax, (0.55, y), 0.72, 0.36, [ROYAL, EMERALD, SLATE][i % 3], text=f"L{i}", fontsize=7.8)
        ax.text(1.42, y + 0.18, f"tokens {i*16:02d}-{i*16+15:02d}", va="center", fontsize=7.5, color=SLATE)
        for s in range(4):
            add_rect(ax, (2.62 + s * 0.18, y + 0.06), 0.14, 0.14, [ROYAL, EMERALD, SLATE][i % 3], alpha=0.65, lw=0.2)
        ax.text(3.42, y + 0.17, "...", fontsize=7.2, color=SLATE)
    ax.text(0.55, 2.58, "Logical view is contiguous; only block ids are exposed to the kernel.", fontsize=7.9, color=SLATE)

    ax.text(5.6, 7.1, "Block Table metadata", fontsize=11.5, weight="bold", color=TEXT, ha="center")
    add_rect(ax, (4.25, 2.22), 2.95, 4.42, "white", ec=SLATE, lw=1.2)
    add_rect(ax, (4.25, 6.24), 2.95, 0.4, SLATE, text="L | P | ref | valid", fontsize=7.6)
    refs = [1, 1, 2, 1, 1, 1]
    for i, p in enumerate(phys):
        y = 5.86 - i * 0.56
        add_rect(ax, (4.42, y), 0.43, 0.33, "#EEF2F6", text=f"L{i}", fontsize=7.1, ec=GRID, text_color=TEXT)
        add_rect(ax, (5.05, y), 0.58, 0.33, [ROYAL, EMERALD, SLATE][i % 3], text=f"P{p}", fontsize=7.1)
        ax.text(5.95, y + 0.17, str(refs[i]), ha="center", va="center", fontsize=7.3, color=TEXT)
        ax.text(6.62, y + 0.17, "yes", ha="center", va="center", fontsize=7.0, color=EMERALD)
        ax.annotate("", xy=(4.25, y + 0.16), xytext=(3.65, logical[i][1] + 0.18), arrowprops=dict(arrowstyle="->", color=[ROYAL, EMERALD, SLATE][i % 3], lw=0.8, alpha=0.72))

    ax.text(10.15, 7.1, "GPU KV memory pool (non-contiguous)", fontsize=11.5, weight="bold", color=TEXT, ha="center")
    physical_positions = {}
    pid_grid = [2, 18, 7, 33, 5, 11, 24, 4, 14, 31, 9, 40, 45, 28, 6, 22]
    used = set(phys)
    for r in range(4):
        for c in range(4):
            idx = r * 4 + c
            pid = pid_grid[idx]
            x = 8.3 + c * 0.88
            y = 5.9 - r * 0.66
            physical_positions[pid] = (x, y)
            is_used = pid in used
            color = [ROYAL, EMERALD, SLATE][phys.index(pid) % 3] if is_used else "#EEF2F6"
            add_rect(ax, (x, y), 0.74, 0.48, color, text=f"P{pid}" if is_used else "free", fontsize=7.1, ec=GRID if not is_used else "white", text_color=TEXT if not is_used else "white")
            ax.text(x + 0.37, y - 0.09, f"0x{0xC000 + pid * 0x200:X}", ha="center", fontsize=5.7, color=SLATE)

    for i, pid in enumerate(phys):
        src_y = 5.86 - i * 0.56 + 0.16
        px, py = physical_positions[pid]
        ax.annotate("", xy=(px, py + 0.24), xytext=(7.2, src_y), arrowprops=dict(arrowstyle="->", color=[ROYAL, EMERALD, SLATE][i % 3], lw=0.9, alpha=0.75))

    ax.text(10.15, 2.72, "Kernel gather order", fontsize=11.2, weight="bold", color=TEXT, ha="center")
    gather_x = 7.85
    for i, pid in enumerate(phys):
        x = gather_x + i * 0.78
        add_rect(ax, (x, 1.92), 0.62, 0.42, [ROYAL, EMERALD, SLATE][i % 3], text=f"P{pid}", fontsize=6.9)
        ax.text(x + 0.31, 1.72, f"K/V{i}", ha="center", fontsize=6.5, color=SLATE)
        if i < len(phys) - 1:
            ax.annotate("", xy=(x + 0.77, 2.13), xytext=(x + 0.62, 2.13), arrowprops=dict(arrowstyle="->", lw=0.75, color=SLATE))
    add_rect(ax, (12.75, 1.75), 0.85, 0.68, EMERALD, text="softmax\nQK^T", fontsize=6.7)
    ax.annotate("", xy=(12.75, 2.1), xytext=(12.33, 2.1), arrowprops=dict(arrowstyle="->", lw=1.0, color=SLATE))

    ax.text(
        4.15,
        0.95,
        "for block_id in logical_blocks:  physical = block_table[block_id];  gather(K_cache[physical], V_cache[physical])",
        ha="center",
        fontsize=8.0,
        color=TEXT,
        bbox=dict(boxstyle="round,pad=0.32", fc=LIGHT_BG, ec=GRID),
    )
    ax.text(10.1, 0.95, "Result: scattered physical pages are reconstructed as one logically continuous attention context.", fontsize=8.0, color=SLATE)
    return save_fig(fig, "slide_05_paged_attention_lookup", tight=False)


def build_slide_6_figure() -> Path:
    fig, ax = plt.subplots(figsize=(10.7, 6.4))
    concurrency = np.array([1, 2, 4, 8, 16, 32, 64])
    x = np.arange(len(concurrency))
    width = 0.25
    hf = np.array([2.0, 3.9, 7.5, 10.5, 0, 0, 0])
    ft = np.array([4.2, 8.0, 15.5, 28.0, 39.0, 0, 0])
    vllm = np.array([21, 42, 82, 158, 300, 555, 930])

    bars_hf = ax.bar(x - width, hf, width, label="HF Transformers", color=SLATE)
    bars_ft = ax.bar(x, ft, width, label="FasterTransformer", color=EMERALD)
    bars_v = ax.bar(x + width, vllm, width, label="vLLM", color=ROYAL)
    ax.set_xlabel("Concurrency (number of active requests)")
    ax.set_ylabel("Throughput (requests/s)")
    ax.set_xticks(x)
    ax.set_xticklabels(concurrency)
    ax.set_ylim(0, 1050)
    figure_title(ax, "Throughput and Latency under Increasing Concurrency", "Representative trend consistent with reported 10-24x vLLM throughput gain.")
    ax.grid(axis="y", color=GRID)

    for bars in [bars_hf, bars_ft, bars_v]:
        for bar in bars:
            h = bar.get_height()
            if h > 0:
                ax.text(bar.get_x() + bar.get_width() / 2, h + 18, f"{h:.0f}", ha="center", va="bottom", fontsize=7.5, color=TEXT)

    ax.text(x[4] - width, 90, "OOM", ha="center", color="#A23B3B", fontsize=8.5, weight="bold", rotation=90)
    ax.text(x[5], 90, "OOM", ha="center", color="#A23B3B", fontsize=8.5, weight="bold", rotation=90)
    ax.text(x[6], 90, "OOM", ha="center", color="#A23B3B", fontsize=8.5, weight="bold", rotation=90)

    ax2 = ax.twinx()
    lat_hf = np.array([145, 168, 230, 410, np.nan, np.nan, np.nan])
    lat_ft = np.array([120, 138, 172, 240, 360, np.nan, np.nan])
    lat_v = np.array([112, 118, 126, 141, 162, 188, 225])
    ax2.plot(x, lat_hf, color=SLATE, marker="o", lw=2, linestyle="--", label="HF latency")
    ax2.plot(x, lat_ft, color=EMERALD, marker="s", lw=2, linestyle="--", label="FT latency")
    ax2.plot(x, lat_v, color=ROYAL, marker="D", lw=2.4, label="vLLM latency")
    ax2.set_ylabel("Average latency (ms/token)")
    ax2.set_ylim(80, 470)
    for xi, val in zip(x, lat_v):
        ax2.text(xi, val + 10, f"{val:.0f}", fontsize=7.2, color=ROYAL, ha="center")

    h1, l1 = ax.get_legend_handles_labels()
    h2, l2 = ax2.get_legend_handles_labels()
    ax.legend(h1 + h2, l1 + l2, loc="upper left", ncol=2, fontsize=8.3)
    return save_fig(fig, "slide_06_throughput_latency")


def build_slide_7_figure() -> Path:
    fig, axes = plt.subplots(1, 2, figsize=(10.7, 5.95), sharey=True)
    seq = ["1k", "2k", "4k", "8k"]
    conc = ["8", "16", "32", "64"]
    vllm = np.array(
        [
            [12, 18, 29, 50],
            [15, 24, 39, 66],
            [22, 36, 58, 92],
            [36, 60, 98, 158],
        ]
    )
    traditional = np.array(
        [
            [21, 39, 76, 151],
            [36, 70, 138, 274],
            [69, 135, 270, 540],
            [132, 264, 528, 1056],
        ]
    )

    mask_trad = traditional > 80
    sns.heatmap(
        np.where(mask_trad, np.nan, traditional),
        ax=axes[0],
        cmap=sns.light_palette(SLATE, as_cmap=True),
        vmin=0,
        vmax=160,
        cbar=False,
        linewidths=0.8,
        linecolor="white",
        annot=False,
    )
    sns.heatmap(
        vllm,
        ax=axes[1],
        cmap=sns.light_palette(ROYAL, as_cmap=True),
        vmin=0,
        vmax=160,
        cbar_kws={"label": "GPU memory footprint (GB)"},
        linewidths=0.8,
        linecolor="white",
        annot=False,
    )

    for ax, title in zip(axes, ["Traditional pre-allocation", "vLLM paged KV blocks"]):
        ax.set_title(title, fontsize=12.5, weight="bold", color=TEXT)
        ax.set_xlabel("Concurrency")
        ax.set_xticklabels(conc)
        ax.set_ylabel("Sequence length" if ax is axes[0] else "")
        ax.set_yticklabels(seq, rotation=0)

    for i in range(traditional.shape[0]):
        for j in range(traditional.shape[1]):
            if mask_trad[i, j]:
                axes[0].add_patch(
                    patches.Rectangle((j, i), 1, 1, facecolor=PALE_RED, edgecolor="white", hatch="///", linewidth=0.8)
                )
                axes[0].text(j + 0.5, i + 0.5, "OOM", ha="center", va="center", fontsize=9, color="#8A2C2C", weight="bold")
            else:
                axes[0].text(j + 0.5, i + 0.5, f"{traditional[i,j]}G", ha="center", va="center", fontsize=8.3, color=TEXT)
            axes[1].text(j + 0.5, i + 0.5, f"{vllm[i,j]}G", ha="center", va="center", fontsize=8.3, color=TEXT)

    fig.suptitle("Memory Stability for Long-context Inference", x=0.02, ha="left", fontsize=15, weight="bold", color=TEXT)
    fig.text(0.02, 0.91, "80GB A100 capacity boundary; vLLM grows linearly with low fragmentation overhead.", fontsize=10, color=SLATE)
    fig.tight_layout(rect=(0, 0, 1, 0.88))
    return save_fig(fig, "slide_07_memory_stability")


def build_slide_8_figure() -> Path:
    labels = ["Throughput", "Latency", "Memory\nUtilization", "Scalability", "Ease of Use"]
    vllm = np.array([9.6, 8.4, 9.7, 9.3, 8.8])
    hf = np.array([3.2, 6.1, 3.5, 4.0, 8.7])
    ft = np.array([6.6, 7.2, 6.2, 6.7, 5.5])
    angles = np.linspace(0, 2 * np.pi, len(labels), endpoint=False).tolist()
    angles += angles[:1]

    fig = plt.figure(figsize=(9.5, 6.4))
    ax = plt.subplot(111, polar=True)
    ax.set_theta_offset(np.pi / 2)
    ax.set_theta_direction(-1)
    ax.set_ylim(0, 10)
    ax.set_yticks([2, 4, 6, 8, 10])
    ax.set_yticklabels(["2", "4", "6", "8", "10"], fontsize=8, color=SLATE)
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(labels, fontsize=10, color=TEXT)
    ax.grid(color=GRID, linewidth=0.9)

    for values, color, label, alpha in [
        (vllm, ROYAL, "vLLM", 0.26),
        (hf, SLATE, "HF Transformers", 0.13),
        (ft, EMERALD, "FasterTransformer", 0.16),
    ]:
        vals = values.tolist() + values[:1].tolist()
        ax.plot(angles, vals, color=color, lw=2.4, label=label)
        ax.fill(angles, vals, color=color, alpha=alpha)
        for angle, val in zip(angles[:-1], values):
            ax.text(angle, val + 0.35, f"{val:.1f}", fontsize=7.6, color=color, ha="center", va="center")

    ax.set_title("Engineering Impact Radar", y=1.12, fontsize=15, weight="bold", color=TEXT)
    ax.legend(loc="lower center", bbox_to_anchor=(0.5, -0.16), ncol=3, fontsize=9)
    return save_fig(fig, "slide_08_radar")


def generate_figures() -> list[Path]:
    return [
        build_slide_1_figure(),
        build_slide_2_figure(),
        build_slide_3_figure(),
        build_slide_4_figure(),
        build_slide_5_figure(),
        build_slide_6_figure(),
        build_slide_7_figure(),
        build_slide_8_figure(),
    ]


def set_run(run, size=14, bold=False, color=RGB["text"]):
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_header(slide, number: int, title: str) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.18))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGB["royal"]
    bar.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.35), Inches(0.28), Inches(12.7), Inches(0.48))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = f"{number:02d}  {title}"
    set_run(r, size=18, bold=True, color=RGB["text"])

    rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.35), Inches(0.82), Inches(12.65), Inches(0.02))
    rule.fill.solid()
    rule.fill.fore_color.rgb = RGBColor(0xD7, 0xDE, 0xE8)
    rule.line.fill.background()


def add_left_image(slide, image_path: Path) -> None:
    # Left pane occupies roughly 65% of the slide width; preserve the figure aspect ratio.
    box_x, box_y, box_w, box_h = Inches(0.33), Inches(0.96), Inches(8.35), Inches(6.12)
    with Image.open(image_path) as img:
        img_w, img_h = img.size
    image_ratio = img_w / img_h
    box_ratio = box_w / box_h
    if image_ratio >= box_ratio:
        width = box_w
        height = int(box_w / image_ratio)
        left = box_x
        top = int(box_y + (box_h - height) / 2)
    else:
        height = box_h
        width = int(box_h * image_ratio)
        left = int(box_x + (box_w - width) / 2)
        top = box_y
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def add_right_panel(slide, heading: str, sections: list[tuple[str, str]]) -> None:
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(8.92),
        Inches(0.96),
        Inches(4.08),
        Inches(6.12),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGB["light"]
    panel.line.color.rgb = RGBColor(0xD7, 0xDE, 0xE8)
    panel.line.width = Pt(1.0)

    tx = slide.shapes.add_textbox(Inches(9.12), Inches(1.13), Inches(3.68), Inches(5.78))
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = heading
    set_run(r, size=15.2, bold=True, color=RGB["royal"])
    p.space_after = Pt(8)

    for label, body in sections:
        p = tf.add_paragraph()
        p.level = 0
        p.space_before = Pt(4)
        p.space_after = Pt(2)
        p.line_spacing = 1.02
        p.text = ""
        r1 = p.add_run()
        r1.text = f"{label}："
        set_run(r1, size=10.6, bold=True, color=RGB["text"])
        r2 = p.add_run()
        r2.text = body
        set_run(r2, size=10.3, bold=False, color=RGB["muted"])


def add_footer(slide) -> None:
    tx = slide.shapes.add_textbox(Inches(0.38), Inches(7.15), Inches(12.55), Inches(0.18))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = "vLLM inference pipeline and performance visualization | Academic schematic, normalized data"
    set_run(r, size=7.8, bold=False, color=RGB["muted"])


def add_title_slide(slide, image_path: Path) -> None:
    add_left_image(slide, image_path)
    add_right_panel(
        slide,
        "vLLM 框架推理流程及其性能增益可视化分析",
        [
            ("主题", "vLLM: 针对大模型推理的高性能显存管理与调度分析。"),
            ("核心摘要", "本报告聚焦 PagedAttention 技术，解析其如何通过非连续显存分配突破 LLM 推理吞吐量瓶颈。"),
            ("微观视角", "围绕 Prefill、Decoding 与 Attention 查找路径，展示 KV Cache 张量从逻辑序列到物理块的演变过程。"),
        ],
    )


def build_presentation(figures: list[Path]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slides = [
        (
            "标题页",
            "vLLM: 高性能显存管理与调度分析",
            [
                ("主题", "vLLM: 针对大模型推理的高性能显存管理与调度分析。"),
                ("核心摘要", "本报告聚焦于 PagedAttention 技术，深入解析其如何通过非连续显存分配突破 LLM 推理的吞吐量瓶颈，并对推理全流程中的张量演变进行微观可视化。"),
            ],
        ),
        (
            "内存痛点：KV Cache 碎片化微观分析",
            "KV Cache 碎片化",
            [
                ("现状", "LLM 生成长度不可预知，传统框架通常为每个请求分配最大长度的连续显存。"),
                ("瓶颈", "静态分配导致内部碎片、外部碎片与预留浪费并存，直接压缩可承载 Batch Size。"),
                ("公式", "Waste = sum_i(max_len - len_i) * 2 * n_layers * H_kv * d_head * bytes。"),
                ("对比", "传统显存利用率常落在 20%-40%，vLLM 通过块式复用可提升至 >96%。"),
            ],
        ),
        (
            "推理流程 (I)：Prefill 阶段的张量与块分配",
            "Prefill: 张量到块",
            [
                ("输入转换", "用户 Prompt 经过一次性前向计算，产生初始 Q/K/V 张量与可复用 KV Cache。"),
                ("分块逻辑", "vLLM 根据逻辑索引 floor(token_id / block_size) 将连续序列切片，典型 Block Size=16。"),
                ("维度演变", "X_prompt: [B,S_prompt,d_model] -> K/V: [L,H_kv,S_prompt,d_head] -> physical blocks P_i。"),
                ("状态", "显存池按需分块，避免传统方式下一次性全量预留最大序列长度。"),
            ],
        ),
        (
            "推理流程 (II)：Decoding 阶段的 Slot 填充细节",
            "Decoding: Slot 级写入",
            [
                ("步进机制", "每次迭代仅产生一个新 Token 及其对应 K_t/V_t。"),
                ("精准插入", "vLLM 无需重排历史 KV，只通过 Block Table 找到当前未满物理块并写入指定 Slot。"),
                ("地址偏移", "slot_id = t mod B；addr = base(P_b) + slot_id * bytes(KV_token)。"),
                ("维度演变", "KV Cache 逻辑长度从 T 增加到 T+1，物理上仅增加一个 Slot 的存储开销。"),
            ],
        ),
        (
            "推理流程 (III)：PagedAttention 查找逻辑",
            "PagedAttention 查表访问",
            [
                ("内存访问", "Attention 计算不再假设连续物理地址，而是通过 Block Table 将逻辑块跳转到物理块。"),
                ("技术优势", "该机制类似操作系统虚拟内存，在物理显存碎片化时仍维持逻辑序列连续。"),
                ("工程含义", "请求可共享、复制或释放块，调度器可以用更细粒度管理正在生成的序列。"),
            ],
        ),
        (
            "性能增益：吞吐量与并发度的阶梯增长",
            "吞吐量与并发",
            [
                ("吞吐量增益", "公开报告中 vLLM 相比 HuggingFace 通常实现 10-24 倍吞吐量提升。"),
                ("Batching 策略", "Continuous Batching 允许请求完成后立即释放块并插入新请求，提高 GPU 活跃度。"),
                ("OOM 边界", "HF 在并发上升时更早触及显存上限，vLLM 因碎片化低而保持近线性扩展。"),
            ],
        ),
        (
            "性能增益：长序列下的内存稳定性",
            "长序列显存稳定性",
            [
                ("显存效率", "通过极低碎片化率（<4%），vLLM 在 80GB A100 上可容纳比基准更多的并发任务。"),
                ("线性增长", "显存占用主要随真实 Token 数增长，而非随最大预留长度跳变。"),
                ("场景价值", "长文本、多轮对话和检索增强生成场景中，单 Token 推理成本下降更明显。"),
            ],
        ),
        (
            "结论与工程意义",
            "结论与展望",
            [
                ("技术总结", "vLLM 通过 PagedAttention 将 KV Cache 管理从连续大块分配转为页式块管理，缓解 LLM 推理的内存墙。"),
                ("工程意义", "更高 Batch Size、更稳定并发和更低碎片化共同转化为吞吐量、成本与服务弹性的提升。"),
                ("未来展望", "面向多模态 VLM 与更大规模 Speculative Decoding，页式显存管理仍是核心基础设施。"),
            ],
        ),
    ]

    for idx, (header, right_heading, sections) in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank)
        add_header(slide, idx, header)
        add_left_image(slide, figures[idx - 1])
        add_right_panel(slide, right_heading, sections)
        add_footer(slide)

    prs.save(OUT)


def main() -> None:
    figures = generate_figures()
    build_presentation(figures)
    print(f"Generated {OUT} with {len(figures)} visualization assets.")


if __name__ == "__main__":
    main()
