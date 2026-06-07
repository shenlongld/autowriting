from __future__ import annotations

import math
from pathlib import Path
from textwrap import wrap

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path("generated")
ASSET_DIR = OUT_DIR / "assets"
OUT_FILE = OUT_DIR / "vllm_sosp2023_pagedattention_analysis.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_CN = "Microsoft YaHei"

NAVY = RGBColor(15, 23, 42)
BLUE = RGBColor(37, 99, 235)
BLUE_DARK = RGBColor(30, 64, 175)
TEAL = RGBColor(13, 148, 136)
ORANGE = RGBColor(245, 158, 11)
RED = RGBColor(220, 38, 38)
GREEN = RGBColor(22, 163, 74)
GRAY_700 = RGBColor(55, 65, 81)
GRAY_500 = RGBColor(107, 114, 128)
GRAY_300 = RGBColor(209, 213, 219)
GRAY_100 = RGBColor(243, 244, 246)
WHITE = RGBColor(255, 255, 255)

P = {
    "navy": (15, 23, 42),
    "blue": (37, 99, 235),
    "blue_dark": (30, 64, 175),
    "blue_light": (219, 234, 254),
    "teal": (13, 148, 136),
    "teal_light": (204, 251, 241),
    "orange": (245, 158, 11),
    "orange_light": (254, 243, 199),
    "red": (220, 38, 38),
    "red_light": (254, 226, 226),
    "green": (22, 163, 74),
    "green_light": (220, 252, 231),
    "slate": (71, 85, 105),
    "gray": (148, 163, 184),
    "grid": (226, 232, 240),
    "panel": (248, 250, 252),
    "white": (255, 255, 255),
}


def emu(x: float):
    return Inches(x)


def font_path(bold: bool = False) -> str:
    candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation2/LiberationSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf",
    ]
    for candidate in candidates:
        if Path(candidate).exists():
            return candidate
    return candidates[0]


def pil_font(size: int, bold: bool = False):
    return ImageFont.truetype(font_path(bold), size)


def text_bbox(draw: ImageDraw.ImageDraw, xy, text: str, font):
    return draw.textbbox(xy, text, font=font)


def centered_text(draw, box, text, font, fill=P["navy"]):
    x1, y1, x2, y2 = box
    bbox = text_bbox(draw, (0, 0), text, font)
    tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
    draw.text((x1 + (x2 - x1 - tw) / 2, y1 + (y2 - y1 - th) / 2 - 2), text, font=font, fill=fill)


def rounded_box(draw, box, fill, outline=None, radius=26, width=3, shadow=True):
    x1, y1, x2, y2 = box
    if shadow:
        draw.rounded_rectangle((x1 + 10, y1 + 12, x2 + 10, y2 + 12), radius=radius, fill=(226, 232, 240))
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=width if outline else 1)


def wrap_text(draw, text: str, font, max_width: int):
    words = text.split()
    lines = []
    line = ""
    for word in words:
        candidate = word if not line else f"{line} {word}"
        bbox = text_bbox(draw, (0, 0), candidate, font)
        if bbox[2] - bbox[0] <= max_width:
            line = candidate
        else:
            if line:
                lines.append(line)
            line = word
    if line:
        lines.append(line)
    return lines


def draw_wrapped(draw, xy, text, font, fill, max_width, line_gap=8):
    x, y = xy
    for line in wrap_text(draw, text, font, max_width):
        draw.text((x, y), line, font=font, fill=fill)
        y += font.size + line_gap
    return y


def arrow(draw, start, end, fill=P["slate"], width=7, head=18):
    x1, y1 = start
    x2, y2 = end
    draw.line((x1, y1, x2, y2), fill=fill, width=width)
    angle = math.atan2(y2 - y1, x2 - x1)
    left = (x2 - head * math.cos(angle - math.pi / 6), y2 - head * math.sin(angle - math.pi / 6))
    right = (x2 - head * math.cos(angle + math.pi / 6), y2 - head * math.sin(angle + math.pi / 6))
    draw.polygon([(x2, y2), left, right], fill=fill)


def base_canvas(title: str, subtitle: str = ""):
    img = Image.new("RGB", (1600, 1320), P["white"])
    draw = ImageDraw.Draw(img)
    draw.rounded_rectangle((28, 28, 1572, 1292), radius=42, fill=P["panel"], outline=P["grid"], width=3)
    draw.text((86, 72), title, font=pil_font(48, True), fill=P["navy"])
    if subtitle:
        draw.text((90, 135), subtitle, font=pil_font(25), fill=P["slate"])
    return img, draw


def metric_badge(draw, x, y, label, value, color):
    rounded_box(draw, (x, y, x + 310, y + 120), fill=P["white"], outline=color, radius=28, width=4)
    draw.text((x + 28, y + 24), value, font=pil_font(36, True), fill=color)
    draw.text((x + 30, y + 76), label, font=pil_font(20), fill=P["slate"])


def save(img: Image.Image, name: str) -> Path:
    path = ASSET_DIR / name
    img.save(path, quality=95)
    return path


def slide01_cover():
    img, draw = base_canvas("vLLM Serving System", "PagedAttention turns KV cache into paged GPU memory")
    nodes = [
        ("Requests", "Prompt + sampling params", P["blue_light"], P["blue"]),
        ("Scheduler", "continuous batching", P["teal_light"], P["teal"]),
        ("Block Manager", "allocate / free pages", P["orange_light"], P["orange"]),
        ("PagedAttention", "block table lookup", P["green_light"], P["green"]),
    ]
    y = 360
    for i, (name, desc, fill, line) in enumerate(nodes):
        x = 125 + i * 350
        rounded_box(draw, (x, y, x + 270, y + 185), fill=fill, outline=line, radius=34, width=5)
        draw.text((x + 28, y + 42), name, font=pil_font(31, True), fill=P["navy"])
        draw_wrapped(draw, (x + 30, y + 92), desc, pil_font(22), P["slate"], 210, 4)
        if i < len(nodes) - 1:
            arrow(draw, (x + 280, y + 92), (x + 335, y + 92), P["slate"], width=7)
    rounded_box(draw, (520, 680, 1085, 865), fill=P["white"], outline=P["blue"], radius=34, width=5)
    draw.text((565, 725), "Result: more live sequences per GPU", font=pil_font(35, True), fill=P["navy"])
    draw.text((580, 790), "same model weights, much better KV cache economics", font=pil_font(24), fill=P["slate"])
    metric_badge(draw, 145, 980, "Throughput vs HF", "22-28x", P["blue"])
    metric_badge(draw, 645, 980, "Throughput vs FT", "2-3x", P["teal"])
    metric_badge(draw, 1145, 980, "KV waste", "<4%", P["green"])
    return save(img, "slide_01_system_overview.png")


def slide02_memory_bottleneck():
    img, draw = base_canvas("KV Cache Memory Waste", "Why contiguous reservation limits batch size")
    chart = (180, 265, 1020, 1040)
    x0, y0, x1, y1 = chart
    draw.line((x0, y1, x1, y1), fill=P["gray"], width=3)
    draw.line((x0, y0, x0, y1), fill=P["gray"], width=3)
    for pct in range(0, 101, 20):
        y = y1 - int((y1 - y0) * pct / 100)
        draw.line((x0 - 12, y, x1, y), fill=P["grid"], width=2)
        draw.text((x0 - 78, y - 14), f"{pct}%", font=pil_font(20), fill=P["slate"])

    stacks = [
        ("Contiguous\nallocation", [("Useful KV", 25, P["blue"]), ("Internal frag.", 25, P["orange"]),
                                    ("External frag.", 20, P["red"]), ("Reservation", 30, P["gray"])]),
        ("PagedAttention", [("Useful KV", 96, P["teal"]), ("Waste", 4, P["orange"])]),
    ]
    bar_w = 190
    for idx, (label, values) in enumerate(stacks):
        bx = x0 + 190 + idx * 390
        current = y1
        for name, pct, color in values:
            h = int((y1 - y0) * pct / 100)
            current -= h
            draw.rounded_rectangle((bx, current, bx + bar_w, current + h), radius=8, fill=color)
            if pct >= 12:
                centered_text(draw, (bx, current, bx + bar_w, current + h), f"{pct}%", pil_font(24, True), P["white"])
        for j, line in enumerate(label.split("\n")):
            centered_text(draw, (bx - 55, y1 + 28 + j * 32, bx + bar_w + 55, y1 + 62 + j * 32), line, pil_font(24, True), P["navy"])

    legend_y = 345
    for name, color in [("Useful KV", P["blue"]), ("Fragmentation", P["orange"]), ("Reservation", P["gray"]), ("vLLM compact", P["teal"])]:
        draw.rounded_rectangle((1110, legend_y, 1148, legend_y + 38), radius=9, fill=color)
        draw.text((1170, legend_y + 4), name, font=pil_font(24), fill=P["navy"])
        legend_y += 62
    rounded_box(draw, (1085, 710, 1450, 955), fill=P["white"], outline=P["blue"], radius=28, width=4)
    draw.text((1125, 750), "Batch Size is memory-bound", font=pil_font(28, True), fill=P["navy"])
    draw_wrapped(draw, (1125, 805), "KV cache grows with batch, context length and layers. Waste directly reduces the number of concurrent requests.", pil_font(23), P["slate"], 290)
    return save(img, "slide_02_memory_bottleneck.png")


def slide03_pagedattention():
    img, draw = base_canvas("PagedAttention Mapping", "Logical token blocks are mapped to physical KV pages")
    # Logical blocks
    draw.text((115, 280), "Logical blocks of one sequence", font=pil_font(29, True), fill=P["navy"])
    logical = []
    for i in range(6):
        x = 120 + i * 130
        logical.append((x + 58, 435))
        rounded_box(draw, (x, 380, x + 108, y := 490), fill=P["blue_light"], outline=P["blue"], radius=20, width=4, shadow=False)
        centered_text(draw, (x, 380, x + 108, y), f"L{i}", pil_font(28, True), P["blue_dark"])

    # Page table
    rounded_box(draw, (555, 600, 870, 1038), fill=P["white"], outline=P["teal"], radius=28, width=4)
    draw.text((600, 633), "Block table", font=pil_font(31, True), fill=P["navy"])
    mappings = [7, 2, 11, 4, 9, 14]
    for i, p in enumerate(mappings):
        y = 700 + i * 52
        draw.rounded_rectangle((605, y, 705, y + 38), radius=8, fill=P["blue_light"])
        draw.rounded_rectangle((725, y, 825, y + 38), radius=8, fill=P["teal_light"])
        centered_text(draw, (605, y, 705, y + 38), f"L{i}", pil_font(20, True), P["blue_dark"])
        centered_text(draw, (725, y, 825, y + 38), f"P{p}", pil_font(20, True), P["teal"])
        arrow(draw, (706, y + 19), (724, y + 19), P["slate"], width=3, head=8)

    # Physical pool
    draw.text((1015, 280), "Physical KV block pool", font=pil_font(29, True), fill=P["navy"])
    pool_positions = {}
    for i in range(16):
        col = i % 4
        row = i // 4
        x = 1010 + col * 125
        y = 365 + row * 120
        used = i in mappings
        fill = P["teal_light"] if used else P["white"]
        outline = P["teal"] if used else P["grid"]
        pool_positions[i] = (x + 52, y + 50)
        rounded_box(draw, (x, y, x + 104, y + 90), fill=fill, outline=outline, radius=18, width=4 if used else 2, shadow=False)
        centered_text(draw, (x, y, x + 104, y + 90), f"P{i}", pil_font(24, True), P["teal"] if used else P["gray"])

    for i, p in enumerate(mappings):
        if i in (0, 1, 3, 5):
            start = logical[i]
            mid = (725, 620)
            end = pool_positions[p]
            arrow(draw, start, (mid[0] - 80, mid[1] - 35 + i * 45), P["gray"], width=3, head=10)
            arrow(draw, (825, 719 + i * 52), end, P["teal"], width=4, head=12)

    rounded_box(draw, (130, 1090, 1400, 1205), fill=P["white"], outline=P["blue"], radius=26, width=3)
    draw.text((175, 1125), "Key effect: no requirement for contiguous KV memory; only the last block can be partially unused.", font=pil_font(29, True), fill=P["navy"])
    return save(img, "slide_03_pagedattention_mapping.png")


def slide04_pipeline():
    img, draw = base_canvas("vLLM Inference Loop", "Scheduler, block manager and kernels form a tight decode loop")
    steps = [
        ("Request\nQueue", "arrival + prompt"),
        ("Scheduler", "select runnable seqs"),
        ("Block\nManager", "page allocation"),
        ("PagedAttention\nKernel", "read KV pages"),
        ("Token\nOutput", "append + stream"),
    ]
    colors = [(P["blue_light"], P["blue"]), (P["teal_light"], P["teal"]), (P["orange_light"], P["orange"]), (P["green_light"], P["green"]), (P["blue_light"], P["blue"])]
    y = 380
    centers = []
    for i, ((title, desc), (fill, line)) in enumerate(zip(steps, colors)):
        x = 95 + i * 295
        rounded_box(draw, (x, y, x + 220, y + 170), fill=fill, outline=line, radius=30, width=5)
        for j, t in enumerate(title.split("\n")):
            centered_text(draw, (x + 18, y + 33 + j * 31, x + 202, y + 72 + j * 31), t, pil_font(27, True), P["navy"])
        centered_text(draw, (x + 18, y + 120, x + 202, y + 150), desc, pil_font(19), P["slate"])
        centers.append((x + 220, y + 85))
        if i < len(steps) - 1:
            arrow(draw, (x + 226, y + 85), (x + 286, y + 85), P["slate"], width=7)

    # Feedback loop
    draw.line((1390, 585, 1390, 840, 190, 840, 190, 590), fill=P["gray"], width=5)
    arrow(draw, (190, 590), (190, 556), P["gray"], width=5, head=14)
    draw.text((520, 874), "decode step repeats until EOS / max tokens", font=pil_font(28, True), fill=P["slate"])

    cards = [
        ("Continuous batching", "keeps GPU occupied as requests finish at different times", P["blue"]),
        ("Page-level KV reuse", "enables larger live batch under same memory budget", P["teal"]),
        ("Kernel-aware layout", "attention kernel consumes block table directly", P["green"]),
    ]
    for i, (title, body, color) in enumerate(cards):
        x = 130 + i * 470
        rounded_box(draw, (x, 1000, x + 380, 1175), fill=P["white"], outline=color, radius=26, width=4)
        draw.text((x + 28, 1030), title, font=pil_font(27, True), fill=color)
        draw_wrapped(draw, (x + 30, 1083), body, pil_font(21), P["slate"], 315, 4)
    return save(img, "slide_04_inference_loop.png")


def chart_axes(draw, box, ymax, y_ticks, xlabel="", ylabel=""):
    x0, y0, x1, y1 = box
    draw.line((x0, y1, x1, y1), fill=P["gray"], width=3)
    draw.line((x0, y0, x0, y1), fill=P["gray"], width=3)
    for tick in y_ticks:
        y = y1 - int((y1 - y0) * tick / ymax)
        draw.line((x0, y, x1, y), fill=P["grid"], width=2)
        draw.text((x0 - 75, y - 15), str(tick), font=pil_font(20), fill=P["slate"])
    if xlabel:
        centered_text(draw, (x0, y1 + 80, x1, y1 + 120), xlabel, pil_font(23), P["slate"])
    if ylabel:
        draw.text((x0 - 70, y0 - 55), ylabel, font=pil_font(23), fill=P["slate"])


def slide05_throughput():
    img, draw = base_canvas("Throughput Gain", "Illustrative req/min values scaled to official ratios")
    box = (195, 290, 1245, 1030)
    chart_axes(draw, box, ymax=280, y_ticks=[0, 70, 140, 210, 280], xlabel="Models", ylabel="req/min")
    data = {
        "Llama-7B": [10, 100, 260],
        "OPT-13B": [8, 78, 205],
    }
    systems = [("HF", P["gray"]), ("FT", P["orange"]), ("vLLM", P["blue"])]
    x0, y0, x1, y1 = box
    group_w = 440
    bar_w = 72
    for g, (model, values) in enumerate(data.items()):
        base_x = x0 + 175 + g * group_w
        for i, (system, color) in enumerate(systems):
            value = values[i]
            h = int((y1 - y0) * value / 280)
            x = base_x + i * 95
            draw.rounded_rectangle((x, y1 - h, x + bar_w, y1), radius=14, fill=color)
            centered_text(draw, (x - 12, y1 - h - 43, x + bar_w + 12, y1 - h - 4), str(value), pil_font(22, True), color)
        centered_text(draw, (base_x - 25, y1 + 28, base_x + 280, y1 + 65), model, pil_font(25, True), P["navy"])

    for i, (system, color) in enumerate(systems):
        x = 1040 + i * 145
        draw.rounded_rectangle((x, 205, x + 34, 239), radius=8, fill=color)
        draw.text((x + 45, 205), system, font=pil_font(24, True), fill=P["navy"])

    rounded_box(draw, (1075, 690, 1470, 940), fill=P["white"], outline=P["blue"], radius=28, width=4)
    draw.text((1115, 730), "Official takeaway", font=pil_font(29, True), fill=P["blue"])
    draw.text((1120, 790), "vLLM is 22-28x over HF", font=pil_font(25, True), fill=P["navy"])
    draw.text((1120, 838), "and 2-3x over FT", font=pil_font(25, True), fill=P["navy"])
    draw.text((1120, 890), "on A100-class serving workloads", font=pil_font(20), fill=P["slate"])
    return save(img, "slide_05_throughput_gain.png")


def slide06_memory_waste():
    img, draw = base_canvas("Memory Waste Ratio", "PagedAttention reduces wasted KV cache from 60-80% to <4%")
    # Gauge bars
    rows = [("Traditional", 70, P["red"], "60-80%"), ("vLLM", 4, P["green"], "<4%")]
    for i, (name, value, color, label) in enumerate(rows):
        y = 395 + i * 245
        draw.text((145, y - 12), name, font=pil_font(33, True), fill=P["navy"])
        draw.rounded_rectangle((410, y, 1270, y + 86), radius=43, fill=P["grid"])
        fill_w = max(70, int(860 * value / 80))
        draw.rounded_rectangle((410, y, 410 + fill_w, y + 86), radius=43, fill=color)
        draw.text((410 + fill_w + 35, y + 20), label, font=pil_font(35, True), fill=color)
        for tick in [0, 20, 40, 60, 80]:
            x = 410 + int(860 * tick / 80)
            draw.line((x, y + 105, x, y + 128), fill=P["gray"], width=2)
            centered_text(draw, (x - 45, y + 134, x + 45, y + 165), f"{tick}%", pil_font(18), P["slate"])

    rounded_box(draw, (190, 890, 645, 1125), fill=P["white"], outline=P["red"], radius=28, width=4)
    draw.text((230, 930), "Before", font=pil_font(32, True), fill=P["red"])
    draw_wrapped(draw, (232, 990), "over-reservation + fragmentation consumes most of the KV budget", pil_font(23), P["slate"], 350)
    rounded_box(draw, (865, 890, 1320, 1125), fill=P["white"], outline=P["green"], radius=28, width=4)
    draw.text((905, 930), "After", font=pil_font(32, True), fill=P["green"])
    draw_wrapped(draw, (907, 990), "only the last page can be partially unused; allocation is demand-driven", pil_font(23), P["slate"], 350)
    arrow(draw, (680, 1010), (825, 1010), P["blue"], width=9, head=28)
    return save(img, "slide_06_memory_waste.png")


def slide07_sequence_length():
    img, draw = base_canvas("Latency vs Sequence Length", "Longer contexts amplify KV cache management differences")
    box = (195, 285, 1280, 1025)
    chart_axes(draw, box, ymax=1200, y_ticks=[0, 300, 600, 900, 1200], xlabel="Sequence length (tokens)", ylabel="Latency (ms)")
    seq = [512, 1024, 2048, 4096, 8192]
    series = {
        "HF": ([55, 100, 210, 500, 1100], P["gray"]),
        "FT": ([35, 62, 130, 310, 720], P["orange"]),
        "vLLM": ([28, 48, 90, 180, 330], P["blue"]),
    }
    x0, y0, x1, y1 = box
    points_by_name = {}
    for i, length in enumerate(seq):
        x = x0 + int((x1 - x0) * i / (len(seq) - 1))
        draw.line((x, y1, x, y1 + 15), fill=P["gray"], width=2)
        centered_text(draw, (x - 70, y1 + 28, x + 70, y1 + 58), str(length), pil_font(19), P["slate"])
    for name, (values, color) in series.items():
        pts = []
        for i, value in enumerate(values):
            x = x0 + int((x1 - x0) * i / (len(seq) - 1))
            y = y1 - int((y1 - y0) * value / 1200)
            pts.append((x, y))
        points_by_name[name] = pts
        for a, b in zip(pts, pts[1:]):
            draw.line((*a, *b), fill=color, width=7)
        for x, y in pts:
            draw.ellipse((x - 12, y - 12, x + 12, y + 12), fill=color, outline=P["white"], width=4)
    for idx, (name, (_, color)) in enumerate(series.items()):
        x = 970 + idx * 150
        draw.line((x, 215, x + 50, 215), fill=color, width=7)
        draw.ellipse((x + 18, 203, x + 42, 227), fill=color, outline=P["white"], width=3)
        draw.text((x + 62, 202), name, font=pil_font(23, True), fill=P["navy"])
    rounded_box(draw, (955, 785, 1475, 960), fill=P["white"], outline=P["blue"], radius=28, width=4)
    draw.text((995, 825), "vLLM stays flatter", font=pil_font(31, True), fill=P["blue"])
    draw_wrapped(draw, (997, 880), "Advantage becomes more visible when sequence length and KV cache footprint grow.", pil_font(23), P["slate"], 420)
    return save(img, "slide_07_sequence_latency.png")


def slide08_value():
    img, draw = base_canvas("Serving Economics", "PagedAttention shifts the bottleneck from memory waste to scheduling policy")
    center = (800, 625)
    # Triad
    items = [
        ("High\nThroughput", "22-28x vs HF", (430, 360), P["blue"], P["blue_light"]),
        ("Low\nLatency", "flatter long-context curve", (965, 360), P["green"], P["green_light"]),
        ("Memory\nEfficient", "KV waste <4%", (700, 840), P["orange"], P["orange_light"]),
    ]
    for title, subtitle, (x, y), color, fill in items:
        draw.line((center[0], center[1], x + 120, y + 82), fill=P["grid"], width=9)
    draw.ellipse((center[0] - 105, center[1] - 105, center[0] + 105, center[1] + 105), fill=P["navy"])
    centered_text(draw, (center[0] - 90, center[1] - 25, center[0] + 90, center[1] + 25), "vLLM", pil_font(42, True), P["white"])
    for title, subtitle, (x, y), color, fill in items:
        rounded_box(draw, (x, y, x + 270, y + 170), fill=fill, outline=color, radius=32, width=5)
        lines = title.split("\n")
        for i, line in enumerate(lines):
            centered_text(draw, (x + 20, y + 30 + i * 34, x + 250, y + 66 + i * 34), line, pil_font(30, True), P["navy"])
        centered_text(draw, (x + 18, y + 120, x + 252, y + 150), subtitle, pil_font(21), P["slate"])

    # Mini comparison matrix
    rounded_box(draw, (185, 1040, 1415, 1195), fill=P["white"], outline=P["grid"], radius=26, width=3)
    headers = ["Metric", "Traditional stack", "vLLM"]
    xs = [230, 600, 1080]
    for i, h in enumerate(headers):
        draw.text((xs[i], 1070), h, font=pil_font(25, True), fill=P["navy"])
    rows = [("Memory waste", "60-80%", "<4%"), ("Concurrency", "fragmentation-limited", "page-limited"), ("Cost/request", "higher", "lower")]
    for r, row in enumerate(rows):
        y = 1115 + r * 30
        draw.text((xs[0], y), row[0], font=pil_font(19), fill=P["slate"])
        draw.text((xs[1], y), row[1], font=pil_font(19), fill=P["red"] if r == 0 else P["slate"])
        draw.text((xs[2], y), row[2], font=pil_font(19, True), fill=P["green"] if r == 0 else P["blue"])
    return save(img, "slide_08_value_summary.png")


def generate_assets():
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    return [
        slide01_cover(),
        slide02_memory_bottleneck(),
        slide03_pagedattention(),
        slide04_pipeline(),
        slide05_throughput(),
        slide06_memory_waste(),
        slide07_sequence_length(),
        slide08_value(),
    ]


def set_run(run, size=16, bold=False, color=NAVY):
    run.font.name = FONT_CN
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(slide, text, x, y, w, h, size=16, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return box


def add_pill(slide, text, x, y, w, color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, emu(x), emu(y), emu(w), emu(0.32))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = emu(0.08)
    tf.margin_right = emu(0.08)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_run(run, size=10.5, bold=True, color=WHITE)


def add_bullets(slide, bullets, x, y, w, h, size=13.3, color=NAVY):
    box = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ""
        p.level = 0
        p.space_after = Pt(7)
        run = p.add_run()
        run.text = f"• {item}"
        set_run(run, size=size, color=color)
    return box


def add_right_panel(slide, page, title, subtitle, takeaway, bullets, note):
    # right background panel
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, emu(7.28), emu(0.72), emu(5.42), emu(6.02))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(248, 250, 252)
    panel.line.color.rgb = GRAY_300
    panel.line.width = Pt(1)

    add_pill(slide, f"Slide {page:02d} · vLLM / SOSP 2023", 7.55, 0.94, 2.15, BLUE)
    add_text(slide, title, 7.55, 1.32, 4.75, 0.56, size=22, bold=True, color=NAVY)
    add_text(slide, subtitle, 7.57, 1.92, 4.65, 0.42, size=11.5, color=GRAY_700)
    add_text(slide, "核心结论", 7.58, 2.52, 1.0, 0.22, size=11.5, bold=True, color=BLUE_DARK)
    add_text(slide, takeaway, 7.58, 2.82, 4.65, 0.80, size=15, bold=True, color=NAVY)
    add_text(slide, "讲解要点", 7.58, 3.86, 1.0, 0.22, size=11.5, bold=True, color=BLUE_DARK)
    add_bullets(slide, bullets, 7.58, 4.16, 4.65, 1.42, size=12.2, color=NAVY)
    add_text(slide, note, 7.58, 5.92, 4.65, 0.52, size=9.6, color=GRAY_500)
    add_text(slide, f"{page}/8", 12.02, 6.95, 0.45, 0.18, size=8.8, color=GRAY_500, align=PP_ALIGN.RIGHT)


def add_visual(slide, image_path: Path):
    # Left-side image embedded as a high-resolution PNG.
    pic = slide.shapes.add_picture(str(image_path), emu(0.42), emu(0.72), width=emu(6.58), height=emu(6.02))
    border = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, emu(0.42), emu(0.72), emu(6.58), emu(6.02))
    border.fill.background()
    border.line.color.rgb = GRAY_300
    border.line.width = Pt(0.9)
    # Keep border behind the picture unsupported by python-pptx; transparent fill makes it a subtle outline overlay.
    return pic


def make_slide(prs, image_path, page, title, subtitle, takeaway, bullets, note):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_visual(slide, image_path)
    add_right_panel(slide, page, title, subtitle, takeaway, bullets, note)
    return slide


def main():
    assets = generate_assets()
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slides = [
        (
            "vLLM 推理流程及其性能增益可视化分析",
            "基于 PagedAttention 的高吞吐量 LLM 服务系统",
            "vLLM 的核心不是改模型结构，而是把 KV Cache 管理做成类似虚拟内存的分页系统。",
            [
                "目标场景：多请求并发、变长输入/输出、A100 80GB 服务环境。",
                "官方结论：吞吐量相对 HF 约 22-28x，相对 FT 约 2-3x。",
                "演讲主线：内存瓶颈 → PagedAttention → 调度闭环 → 性能收益。",
            ],
            "Reference: Kwon et al., SOSP 2023；专业术语保留 English。",
        ),
        (
            "背景：LLM 推理中的内存瓶颈",
            "KV Cache 随 batch size 与 sequence length 增长，直接压缩可并发请求数。",
            "传统连续预分配把大量显存锁在碎片和预留空间中，导致 batch size 很早触顶。",
            [
                "每个请求生成 token 时都追加 Key/Value 张量，显存占用动态增长。",
                "传统系统为最坏输出长度预留连续空间，造成 internal/external fragmentation。",
                "论文披露传统 KV Cache 管理浪费约 60%-80%，是吞吐瓶颈的根因之一。",
            ],
            "左图为堆叠占比示意，用于解释 60%-80% waste 与 <4% waste 的差异。",
        ),
        (
            "核心机制：PagedAttention 原理",
            "逻辑块到物理块的映射，让 KV Cache 不再要求物理连续。",
            "PagedAttention 将每个序列的 KV Cache 切成固定大小 block，用 block table 连接逻辑顺序与物理位置。",
            [
                "Logical blocks 保持模型看到的 token 顺序；physical blocks 可分散在 GPU cache pool。",
                "新增 token 只在需要时申请 block，释放请求时可立即回收物理块。",
                "浪费只出现在序列末尾未填满的 block，因此实测内存浪费低于 4%。",
            ],
            "该机制借鉴 OS virtual memory / paging，但服务对象是 GPU 上的 KV Cache。",
        ),
        (
            "vLLM 推理全流程解析",
            "Scheduler、Block Manager 与 PagedAttention kernel 共同构成 decode loop。",
            "吞吐提升来自系统级闭环：调度层保持 GPU 忙碌，内存层让更多请求能同时留在 batch 中。",
            [
                "Scheduler 按可运行状态组织 continuous batching，应对请求长短不一。",
                "Block Manager 负责 KV block 的分配、回收和 block table 更新。",
                "PagedAttention kernel 根据映射读取非连续 KV block，避免搬运成连续缓存。",
            ],
            "这一页适合用 45 秒讲清楚 vLLM engine 的运行路径。",
        ),
        (
            "性能分析：吞吐量增益",
            "基于官方倍率校准的 req/min 可视化：vLLM 对 HF/FT 均显著领先。",
            "vLLM 在 Llama/OPT 等模型上通过更大的有效 batch 获得 20x+ 对 HF、2x+ 对 FT 的吞吐优势。",
            [
                "HF 基线缺少面向高并发 serving 的细粒度调度与 KV cache 管理。",
                "FT 具备优化 kernel，但在动态请求与 KV cache 碎片场景下仍受限。",
                "vLLM 的优势在长序列、大模型、复杂 decoding 策略下更明显。",
            ],
            "柱状图为倍率关系可视化，不把示意 req/min 当作论文逐点原始表格。",
        ),
        (
            "性能分析：内存浪费率对比",
            "从 60%-80% 到 <4%：显存利用率改善直接转化为可服务并发数。",
            "PagedAttention 的经济性体现在单位 GPU 能容纳更多 live sequences，降低单位请求显存成本。",
            [
                "传统系统的浪费主要来自预留空间和碎片化，而非模型权重本身。",
                "vLLM 的 block 粒度分配使显存容量更接近真实 token 需求。",
                "当请求长度高度不均匀时，分页管理比连续分配更稳健。",
            ],
            "图中 Traditional 取 70% 作为 60%-80% 区间中位示意；vLLM 标注 <4%。",
        ),
        (
            "性能分析：不同序列长度下的表现",
            "长上下文会放大 KV Cache 管理差异，vLLM 的时延增长更平缓。",
            "随着 sequence length 增长，传统系统更快进入内存压力区；vLLM 因分页式 KV 管理保持更稳定的 batch。",
            [
                "短序列时计算可能占主导，系统间差异相对收敛。",
                "长序列时 KV Cache 占用成为主导，碎片/预留浪费会触发排队与 batch 缩小。",
                "论文总结：长序列、大模型、复杂 decoding 下 vLLM 改善更突出。",
            ],
            "折线图为趋势示意，用于解释论文结论而非替代原始实验曲线。",
        ),
        (
            "结论与应用价值",
            "PagedAttention 改变了 LLM serving 的成本结构：同样 GPU，服务更多请求。",
            "vLLM 把不可控的动态 KV Cache 显存问题，重构为可调度、可分页、可回收的系统资源管理问题。",
            [
                "对平台方：提升 GPU utilization，降低单位请求成本。",
                "对业务方：在高并发与长文本场景下获得更稳定 SLA。",
                "对系统设计：LLM 推理优化应同时关注 kernel、scheduler 与 memory manager。",
            ],
            "6 分钟收束：vLLM 的价值不是单点优化，而是内存管理与调度的系统化协同。",
        ),
    ]

    for idx, (text, asset) in enumerate(zip(slides, assets), start=1):
        title, subtitle, takeaway, bullets, note = text
        make_slide(prs, asset, idx, title, subtitle, takeaway, bullets, note)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_FILE)
    print(f"Generated {OUT_FILE} with {len(prs.slides)} slides")
    print(f"Generated {len(assets)} embedded chart images in {ASSET_DIR}")


if __name__ == "__main__":
    main()
